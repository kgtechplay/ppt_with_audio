from __future__ import annotations

from pathlib import Path
import sys

from pptx import Presentation

from template_layout_detector import TemplateLayoutProfile
from template_style_detector import (
    detect_content_slide_style,
    detect_title_slide_style,
    drop_cap_body_pair,
    looks_like_footer_text,
)


# Prints text without crashing when the Windows console cannot encode a symbol.
def _safe_print(text: str) -> None:
    encoding = sys.stdout.encoding or "utf-8"
    safe_text = text.encode(encoding, errors="replace").decode(encoding)
    print(safe_text, flush=True)


# Writes build progress to stdout and an optional debug log file.
class _DebugLogger:
    # Initializes or clears the debug log file.
    def __init__(self, log_path: Path | None):
        self.log_path = log_path
        if self.log_path:
            self.log_path.parent.mkdir(parents=True, exist_ok=True)
            self.log_path.write_text("", encoding="utf-8")

    # Appends a single formatted debug line.
    def write(self, message: str) -> None:
        line = f"[template-com] {message}"
        _safe_print(line)
        if self.log_path:
            with self.log_path.open("a", encoding="utf-8") as file:
                file.write(line + "\n")


# Shortens long text for readable debug messages.
def _preview(text: str, limit: int = 90) -> str:
    text = " ".join((text or "").split())
    if len(text) <= limit:
        return text
    return text[: limit - 3] + "..."


# Formats a shape item dict into a compact debug string.
def _item_summary(item) -> str:
    font_part = f" font={item.get('font_size'):.1f}" if item.get("font_size") else ""
    return (
        f"id={item.get('id')} parents={item.get('parent_ids', ())} "
        f"box=({item.get('left'):.0f},{item.get('top'):.0f},"
        f"{item.get('width'):.0f},{item.get('height'):.0f}){font_part} "
        f"text='{_preview(item.get('text', ''))}'"
    )


PLACEHOLDER_WORDS = (
    "title",
    "subtitle",
    "description",
    "presenter",
    "date",
    "category",
    "section label",
    "key point",
    "highlight",
    "metric",
    "insight",
    "callout",
    "short card description",
    "supporting",
    "goes right here",
    "add a sentence",
)

# Reads a COM shape's position and size.
def _shape_bounds(shape) -> tuple[float, float, float, float] | None:
    try:
        return (
            float(shape.Left),
            float(shape.Top),
            float(shape.Width),
            float(shape.Height),
        )
    except Exception:
        return None


# Reads a COM shape's stable slide-local id.
def _shape_id(shape) -> int | None:
    try:
        return int(shape.Id)
    except Exception:
        return None


# Reads the largest available font size from a COM text shape.
def _shape_font_size(shape) -> float | None:
    try:
        text_range = shape.TextFrame.TextRange
        font_size = float(text_range.Font.Size)
        if font_size > 0:
            return font_size
    except Exception:
        pass

    try:
        text_range = shape.TextFrame.TextRange
        sizes = []
        for index in range(1, text_range.Runs().Count + 1):
            run_size = float(text_range.Runs(index).Font.Size)
            if run_size > 0:
                sizes.append(run_size)
        return max(sizes) if sizes else None
    except Exception:
        return None


# Walks top-level and grouped shapes while preserving parent group ids.
def _iter_shape_tree(shapes_collection, parent_ids=()):
    for index in range(1, shapes_collection.Count + 1):
        shape = shapes_collection(index)
        shape_id = _shape_id(shape)
        if shape_id is None:
            continue

        yield shape, parent_ids

        try:
            group_items = shape.GroupItems
            group_count = group_items.Count
        except Exception:
            group_count = 0

        if group_count:
            yield from _iter_shape_tree(group_items, parent_ids + (shape_id,))


# Extracts stripped text from a COM shape when it has a text frame.
def _text(shape) -> str:
    try:
        if not shape.HasTextFrame:
            return ""
        if not shape.TextFrame.HasText:
            return ""
        return shape.TextFrame.TextRange.Text.strip()
    except Exception:
        return ""


# Collects text-bearing shapes from a slide as detector-friendly dictionaries.
def _text_shapes(slide):
    shapes = []
    for shape, parent_ids in _iter_shape_tree(slide.Shapes):
        text = _text(shape)
        if not text:
            continue
        bounds = _shape_bounds(shape)
        if bounds is None:
            continue
        left, top, width, height = bounds
        shape_id = _shape_id(shape)
        font_size = _shape_font_size(shape)
        shapes.append(
            {
                "shape": shape,
                "id": shape_id,
                "parent_ids": parent_ids,
                "text": text,
                "lower": text.lower(),
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "area": width * height,
                "font_size": font_size,
            }
        )
    return sorted(shapes, key=lambda item: (item["top"], item["left"]))


# Collects all shapes from a slide as geometry dictionaries for deletion/protection logic.
def _all_shape_items(slide):
    items = []
    for shape, parent_ids in _iter_shape_tree(slide.Shapes):
        bounds = _shape_bounds(shape)
        shape_id = _shape_id(shape)
        if bounds is None or shape_id is None:
            continue
        left, top, width, height = bounds
        items.append(
            {
                "shape": shape,
                "id": shape_id,
                "parent_ids": parent_ids,
                "left": left,
                "top": top,
                "width": width,
                "height": height,
                "right": left + width,
                "bottom": top + height,
                "area": width * height,
                "cx": left + width / 2,
                "cy": top + height / 2,
                "text": _text(shape),
            }
        )
    return items


# Looks up a PowerPoint slide by its SlideID.
def _slide_by_id(presentation, slide_id):
    return presentation.Slides.FindBySlideID(slide_id)


# Duplicates a source slide and moves the duplicate to the end of the deck.
def _duplicate_to_end(presentation, source_slide_id):
    source = _slide_by_id(presentation, source_slide_id)
    duplicate_range = source.Duplicate()
    duplicate = duplicate_range.Item(1)
    duplicate.MoveTo(presentation.Slides.Count)
    return presentation.Slides(presentation.Slides.Count)


# Replaces all text in a COM text shape.
def _set_shape_text(shape, text: str) -> None:
    text_range = shape.TextFrame.TextRange
    text_range.Text = text


# Sets a text shape's font size across its text range.
def _set_shape_font_size(shape, font_size: float) -> None:
    try:
        shape.TextFrame.TextRange.Font.Size = font_size
    except Exception:
        pass


# Clears text from a COM text shape without failing the build.
def _clear_shape_text(shape) -> None:
    try:
        shape.TextFrame.TextRange.Text = ""
    except Exception:
        pass


# Deletes a COM shape without failing the build.
def _delete_shape(shape) -> None:
    try:
        shape.Delete()
    except Exception:
        pass


# Computes the slide area used for relative geometry thresholds.
def _slide_area(slide) -> float:
    try:
        setup = slide.Parent.PageSetup
        return float(setup.SlideWidth) * float(setup.SlideHeight)
    except Exception:
        return 720 * 540


# Reads the slide dimensions used for fit calculations.
def _slide_size(slide) -> tuple[float, float]:
    try:
        setup = slide.Parent.PageSetup
        return float(setup.SlideWidth), float(setup.SlideHeight)
    except Exception:
        return 720, 540


# Checks whether a point falls inside a shape item rectangle.
def _contains_point(item, x, y) -> bool:
    return item["left"] <= x <= item["right"] and item["top"] <= y <= item["bottom"]


# Checks whether two shape item rectangles overlap.
def _overlaps(a, b) -> bool:
    return not (
        a["right"] < b["left"]
        or a["left"] > b["right"]
        or a["bottom"] < b["top"]
        or a["top"] > b["bottom"]
    )


# Returns a copy of a shape item rectangle expanded by padding.
def _expanded(item, padding) -> dict:
    return {
        **item,
        "left": item["left"] - padding,
        "top": item["top"] - padding,
        "right": item["right"] + padding,
        "bottom": item["bottom"] + padding,
    }


# Returns the rendered text bounds reported by PowerPoint.
def _text_render_size(shape) -> tuple[float, float] | None:
    try:
        text_range = shape.TextFrame2.TextRange
        return float(text_range.BoundWidth), float(text_range.BoundHeight)
    except Exception:
        pass

    try:
        text_range = shape.TextFrame.TextRange
        return float(text_range.BoundWidth), float(text_range.BoundHeight)
    except Exception:
        return None


# Checks whether rendered text is larger than its shape box.
def _text_overflows_shape(shape, padding=4) -> bool:
    bounds = _shape_bounds(shape)
    render_size = _text_render_size(shape)
    if bounds is None or render_size is None:
        return False

    _, _, width, height = bounds
    rendered_width, rendered_height = render_size
    return rendered_width > width - padding or rendered_height > height - padding


# Finds the rightmost edge a shape can grow to without hitting another shape.
def _available_right_edge(slide, item, padding=10) -> float:
    slide_width, _ = _slide_size(slide)
    limit = slide_width - padding
    item_bottom = item["top"] + item["height"]

    for other in _all_shape_items(slide):
        if other["id"] == item["id"]:
            continue
        overlaps_vertically = not (
            other["bottom"] < item["top"] + padding
            or other["top"] > item_bottom - padding
        )
        if overlaps_vertically and other["left"] >= item["left"] + item["width"]:
            limit = min(limit, other["left"] - padding)

    return max(item["left"] + item["width"], limit)


# Finds the lowest edge a shape can grow to without hitting another shape.
def _available_bottom_edge(slide, item, padding=10) -> float:
    _, slide_height = _slide_size(slide)
    limit = slide_height - padding
    item_right = item["left"] + item["width"]

    for other in _all_shape_items(slide):
        if other["id"] == item["id"]:
            continue
        overlaps_horizontally = not (
            other["right"] < item["left"] + padding
            or other["left"] > item_right - padding
        )
        if overlaps_horizontally and other["top"] >= item["top"] + item["height"]:
            limit = min(limit, other["top"] - padding)

    return max(item["top"] + item["height"], limit)


# Expands a text box into empty space before resorting to font shrinking.
def _grow_text_shape_to_fit(slide, item, logger=None) -> None:
    shape = item["shape"]
    render_size = _text_render_size(shape)
    if render_size is None:
        return

    rendered_width, rendered_height = render_size
    right_edge = _available_right_edge(slide, item)
    max_width = max(item["width"], right_edge - item["left"])
    if rendered_height > item["height"] - 4:
        desired_width = max_width
    else:
        desired_width = min(max_width, max(item["width"], rendered_width + 12))
    if desired_width > item["width"] + 1:
        try:
            shape.Width = desired_width
            if logger:
                logger.write(
                    f"expanded text box width id={item['id']} "
                    f"from {item['width']:.0f} to {desired_width:.0f}"
                )
        except Exception:
            pass

    bounds = _shape_bounds(shape)
    if bounds is None:
        return

    left, top, width, height = bounds
    render_size = _text_render_size(shape)
    rendered_height = render_size[1] if render_size is not None else rendered_height
    current_item = {**item, "left": left, "top": top, "width": width, "height": height}
    bottom_edge = _available_bottom_edge(slide, current_item)
    max_height = max(height, bottom_edge - top)
    desired_height = min(max_height, max(height, rendered_height + 10))
    if desired_height > height + 1:
        try:
            shape.Height = desired_height
            if logger:
                logger.write(
                    f"expanded text box height id={item['id']} "
                    f"from {height:.0f} to {desired_height:.0f}"
                )
        except Exception:
            pass


# Reduces font size until rendered text fits inside its box.
def _shrink_text_to_fit(shape, item, logger=None, min_font_size=8.0) -> None:
    font_size = _shape_font_size(shape)
    if not font_size:
        return

    current_size = font_size
    while current_size > min_font_size and _text_overflows_shape(shape):
        current_size = max(min_font_size, current_size - 1)
        _set_shape_font_size(shape, current_size)

    if logger and current_size < font_size:
        logger.write(
            f"reduced font size id={item['id']} from {font_size:.1f} to {current_size:.1f}"
        )


# Fits generated text shapes by growing available box space, then shrinking text.
def _fit_generated_text_shapes(slide, generated_text_ids, logger=None) -> None:
    if not generated_text_ids:
        return

    for item in _text_shapes(slide):
        if item["id"] not in generated_text_ids:
            continue
        if not item["text"].strip():
            continue
        if not _text_overflows_shape(item["shape"]):
            continue

        if logger:
            logger.write(f"fit generated text shape {_item_summary(item)}")
        _grow_text_shape_to_fit(slide, item, logger=logger)
        _shrink_text_to_fit(item["shape"], item, logger=logger)


# Sets all text shapes in a repeated role to the smallest current font size.
def _normalize_font_size_for_ids(slide, shape_ids, role_name, logger=None) -> None:
    shape_ids = set(shape_ids or [])
    if not shape_ids:
        return

    items = [
        item
        for item in _text_shapes(slide)
        if item["id"] in shape_ids and item["text"].strip() and item.get("font_size")
    ]
    if len(items) < 2:
        return

    target_size = min(item["font_size"] for item in items)
    for item in items:
        if abs(item["font_size"] - target_size) < 0.1:
            continue
        _set_shape_font_size(item["shape"], target_size)
        if logger:
            logger.write(
                f"normalized {role_name} font size id={item['id']} "
                f"from {item['font_size']:.1f} to {target_size:.1f}"
            )


# Deletes an unused text component plus nearby container/decorative shapes.
def _delete_unused_component(slide, text_item, protected_ids=None, logger=None) -> None:
    protected_ids = set(protected_ids or [])
    shape_items = _all_shape_items(slide)
    text_id = int(text_item["id"])
    protected_parent_ids = {
        parent_id
        for item in shape_items
        if item["id"] in protected_ids
        for parent_id in item.get("parent_ids", ())
    }
    text_box = {
        "left": text_item["left"],
        "top": text_item["top"],
        "right": text_item["left"] + text_item["width"],
        "bottom": text_item["top"] + text_item["height"],
        "area": text_item["area"],
        "cx": text_item["left"] + text_item["width"] / 2,
        "cy": text_item["top"] + text_item["height"] / 2,
    }
    slide_area = _slide_area(slide)

    containers = [
        item
        for item in shape_items
        if item["id"] != text_id
        and item["id"] not in protected_parent_ids
        and not item["text"]
        and item["area"] > text_box["area"] * 1.4
        and item["area"] < slide_area * 0.45
        and _contains_point(item, text_box["cx"], text_box["cy"])
    ]

    if containers:
        container = min(containers, key=lambda item: item["area"])
        delete_region = _expanded(container, 14)
        delete_ids = {
            item["id"]
            for item in shape_items
            if item["area"] < slide_area * 0.45
            and item["id"] not in protected_parent_ids
            and (
                _contains_point(delete_region, item["cx"], item["cy"])
                or _overlaps(delete_region, item)
            )
        }
    else:
        container = None
        delete_region = _expanded(text_box, 36)
        delete_ids = {
            item["id"]
            for item in shape_items
            if item["id"] == text_id
            or (
                item["area"] < slide_area * 0.15
                and item["id"] not in protected_parent_ids
                and _contains_point(delete_region, item["cx"], item["cy"])
            )
        }

    delete_ids = delete_ids - protected_ids - protected_parent_ids

    if logger:
        container_text = _item_summary(container) if container else "none"
        logger.write(
            "delete unused component from "
            f"{_item_summary(text_item)} | container={container_text} | "
            f"protected={sorted(protected_ids)} protected_parents={sorted(protected_parent_ids)} "
            f"delete_ids={sorted(delete_ids)}"
        )

    for item in reversed(shape_items):
        if item["id"] in delete_ids:
            _delete_shape(item["shape"])
    return delete_ids


# Finds the smallest non-text container shape around a text component.
def _component_container(slide, text_item):
    shape_items = _all_shape_items(slide)
    text_box = {
        "left": text_item["left"],
        "top": text_item["top"],
        "right": text_item["left"] + text_item["width"],
        "bottom": text_item["top"] + text_item["height"],
        "area": text_item["area"],
        "cx": text_item["left"] + text_item["width"] / 2,
        "cy": text_item["top"] + text_item["height"] / 2,
    }
    slide_area = _slide_area(slide)
    containers = [
        item
        for item in shape_items
        if item["id"] != text_item["id"]
        and not item["text"]
        and item["area"] > text_box["area"] * 1.4
        and item["area"] < slide_area * 0.45
        and _contains_point(item, text_box["cx"], text_box["cy"])
    ]
    return min(containers, key=lambda item: item["area"]) if containers else None


# Detects placeholder/template text that should be removed after replacement.
def _looks_like_template_text(text: str) -> bool:
    lower = text.lower()
    return any(word in lower for word in PLACEHOLDER_WORDS)


# Decides whether an unmapped content text shape should be deleted.
def _should_delete_unmapped_content_text(slide, item) -> bool:
    if looks_like_footer_text(item["text"]):
        return False

    return True


# Formats generated bullet strings for a single body text box.
def _format_body_text(bullets: list[str]) -> str:
    return "\r".join(bullets)


# Writes bullets into a body shape, including special handling for drop-cap layouts.
def _set_body_text(slide, body_shape, bullets: list[str], shapes, protected_ids, logger=None):
    drop_cap_item, drop_cap_partner = drop_cap_body_pair(shapes)
    body_text = _format_body_text(bullets)

    if body_shape is drop_cap_partner and drop_cap_item is not None:
        if bullets:
            first_text = " ".join(bullets[0].split())
            first_letter = first_text[:1]
            remainder = first_text[1:].lstrip()
            remaining_lines = [remainder] if remainder else []
            remaining_lines.extend(bullets[1:])
            _set_shape_text(drop_cap_item["shape"], first_letter)
            _set_shape_text(drop_cap_partner["shape"], "\r".join(remaining_lines))
            protected_ids.add(drop_cap_item["id"])
            protected_ids.add(drop_cap_partner["id"])
            if logger:
                logger.write(
                    "set drop-cap body pair: "
                    f"drop_cap={_item_summary(drop_cap_item)} -> '{_preview(first_letter)}', "
                    f"body={_item_summary(drop_cap_partner)} with {len(remaining_lines)} line(s)"
                )
            return

    _set_shape_text(body_shape["shape"], body_text)
    protected_ids.add(body_shape["id"])
    if logger:
        logger.write(
            f"set body shape id={body_shape['id']} with {len(bullets)} bullet(s): "
            f"{[_preview(bullet, 50) for bullet in bullets]}"
        )


# Protects the full visual card/container region around a repeated component.
def _protect_component_region(slide, component_item, protected_ids):
    container = _component_container(slide, component_item)
    if container is None:
        protected_ids.add(component_item["id"])
        return

    shape_items = _all_shape_items(slide)
    region = _expanded(container, 14)
    for item in shape_items:
        if (
            _contains_point(region, item["cx"], item["cy"])
            or _overlaps(region, item)
        ):
            protected_ids.add(item["id"])


# Finds the description text shape paired with a repeated card component.
def _component_description_item(slide, component_item, all_text_items):
    container = _component_container(slide, component_item)
    if container is None:
        return None
    region = _expanded(container, 8)
    candidates = [
        item
        for item in all_text_items
        if item["id"] != component_item["id"]
        and _contains_point(region, item["left"] + item["width"] / 2, item["top"] + item["height"] / 2)
        and item["top"] >= component_item["top"]
    ]
    return max(candidates, key=lambda item: item["top"]) if candidates else None


# Splits one bullet into card title and card description text.
def _split_component_value(value: str) -> tuple[str, str]:
    value = " ".join(value.split())
    if not value:
        return "", ""
    separators = [": ", " - ", ". "]
    for separator in separators:
        if separator in value:
            head, tail = value.split(separator, 1)
            if head and tail:
                return head.strip(), tail.strip()
    return value, ""


# Fills repeated card components with bullets and cleans up unused cards.
def _fill_repeated_card_components(
    slide,
    components,
    values,
    protected_ids,
    all_text_items,
    logger=None,
    description_by_component_id=None,
):
    used_ids = set()
    deleted_ids = set()
    used_heading_ids = set()
    used_description_ids = set()
    description_by_component_id = description_by_component_id or {}

    for item, value in zip(components, values):
        title_text, description_text = _split_component_value(value)
        description_item = description_by_component_id.get(item["id"])
        if description_item is None:
            description_item = _component_description_item(slide, item, all_text_items)

        _set_shape_text(item["shape"], title_text)
        if logger:
            logger.write(
                f"fill card title {_item_summary(item)} with '{_preview(title_text)}'"
            )
        protected_ids.add(item["id"])
        used_ids.add(item["id"])
        used_heading_ids.add(item["id"])

        if description_item is not None:
            if description_text:
                _set_shape_text(description_item["shape"], description_text)
                if logger:
                    logger.write(
                        "fill card description "
                        f"{_item_summary(description_item)} with '{_preview(description_text)}'"
                    )
            else:
                _clear_shape_text(description_item["shape"])
                if logger:
                    logger.write(
                        f"clear empty card description {_item_summary(description_item)}"
                    )
            protected_ids.add(description_item["id"])
            used_ids.add(description_item["id"])
            used_description_ids.add(description_item["id"])

        _protect_component_region(slide, item, protected_ids)

    for item in components[len(values):]:
        if item["id"] in deleted_ids:
            continue
        deleted_ids.update(_delete_unused_component(slide, item, protected_ids=protected_ids, logger=logger))

    return used_ids, deleted_ids, used_heading_ids, used_description_ids


# Deletes repeated components that are not used by the chosen replacement strategy.
def _delete_repeated_components(slide, components, protected_ids, logger=None):
    deleted_ids = set()
    for item in components:
        if item["id"] in deleted_ids:
            continue
        deleted_ids.update(_delete_unused_component(slide, item, protected_ids=protected_ids, logger=logger))
    return deleted_ids


# Replaces text on the duplicated title sample slide.
def _replace_title_slide_text(slide, title: str, logger=None) -> None:
    shapes = _text_shapes(slide)
    shape_by_id = {item["id"]: item for item in shapes}
    style = detect_title_slide_style(shapes)
    title_shape = shape_by_id.get(style.title_id)
    subtitle_shape = shape_by_id.get(style.subtitle_id)
    protected_ids = set()
    generated_text_ids = set()

    if logger:
        logger.write(f"title slide text shape count={len(shapes)}")
        for item in shapes:
            logger.write(f"title candidate: {_item_summary(item)}")
        logger.write(
            "title mapping: "
            f"title={_item_summary(title_shape) if title_shape else 'none'} | "
            f"subtitle={_item_summary(subtitle_shape) if subtitle_shape else 'none'}"
        )

    for item in shapes:
        if item is title_shape:
            _set_shape_text(item["shape"], title)
            if logger:
                logger.write(f"set title shape id={item['id']} to '{_preview(title)}'")
            protected_ids.add(item["id"])
            generated_text_ids.add(item["id"])
        elif item is subtitle_shape:
            _set_shape_text(item["shape"], "Generated with OpenAI")
            if logger:
                logger.write(f"set subtitle shape id={item['id']} to 'Generated with OpenAI'")
            protected_ids.add(item["id"])
            generated_text_ids.add(item["id"])

    deleted_ids = set()
    for item in shapes:
        if item["id"] in deleted_ids:
            continue
        if item is title_shape or item is subtitle_shape:
            continue
        if looks_like_footer_text(item["text"]):
            if logger:
                logger.write(f"preserve title footer text shape {_item_summary(item)}")
            continue
        if logger:
            logger.write(f"delete unmapped title text shape {_item_summary(item)}")
        _delete_shape(item["shape"])
        deleted_ids.add(item["id"])

    _fit_generated_text_shapes(slide, generated_text_ids, logger=logger)


# Replaces text on a duplicated content sample slide.
def _replace_content_slide_text(slide, title: str, bullets: list[str], logger=None) -> None:
    shapes = _text_shapes(slide)
    shape_by_id = {item["id"]: item for item in shapes}
    style = detect_content_slide_style(shapes)
    title_shape = shape_by_id.get(style.title_id)
    repeated_components = [
        shape_by_id[shape_id]
        for shape_id in style.repeated_component_ids
        if shape_id in shape_by_id
    ]
    body_shape = shape_by_id.get(style.body_id)
    description_by_component_id = {
        component_id: shape_by_id[description_id]
        for component_id, description_id in zip(
            style.repeated_component_ids,
            style.repeated_component_description_ids,
        )
        if component_id in shape_by_id and description_id in shape_by_id
    }
    use_body_region = style.use_body_region
    use_repeated_components = style.use_repeated_components
    active_body_shape = body_shape if use_body_region else None
    protected_ids = set()
    deleted_ids = set()
    generated_text_ids = set()
    repeated_heading_ids = set()
    repeated_description_ids = set()

    if logger:
        logger.write(f"content slide '{_preview(title)}' text shape count={len(shapes)}")
        for item in shapes:
            logger.write(f"content candidate: {_item_summary(item)}")
        logger.write(
            "content mapping: "
            f"title={_item_summary(title_shape) if title_shape else 'none'} | "
            f"body={_item_summary(active_body_shape) if active_body_shape else 'none'} | "
            f"use_body_region={use_body_region} | "
            f"repeated_count={len(repeated_components)} | "
            f"use_repeated_components={use_repeated_components}"
        )
        logger.write(
            "style model: "
            f"drop_cap_id={style.drop_cap_id}, drop_cap_body_id={style.drop_cap_body_id}, "
            f"repeated_description_ids={style.repeated_component_description_ids}, "
            f"deletable_text_ids={sorted(style.deletable_text_ids)}, "
            f"footer_ids={sorted(style.footer_ids)}"
        )
        for item in repeated_components:
            logger.write(f"repeated component candidate: {_item_summary(item)}")

    for item in shapes:
        if item is title_shape:
            _set_shape_text(item["shape"], title)
            if logger:
                logger.write(f"set content title shape id={item['id']} to '{_preview(title)}'")
            protected_ids.add(item["id"])
            generated_text_ids.add(item["id"])
        elif item is active_body_shape:
            _set_body_text(slide, active_body_shape, bullets, shapes, protected_ids, logger=logger)
            generated_text_ids.update(protected_ids)

    filled_component_ids = set()
    if use_repeated_components:
        (
            filled_component_ids,
            component_deleted_ids,
            repeated_heading_ids,
            repeated_description_ids,
        ) = _fill_repeated_card_components(
            slide,
            repeated_components,
            bullets,
            protected_ids,
            shapes,
            logger=logger,
            description_by_component_id=description_by_component_id,
        )
        deleted_ids.update(component_deleted_ids)
        generated_text_ids.update(filled_component_ids)
    elif repeated_components:
        if logger:
            logger.write("body region selected; deleting unused repeated components")
        deleted_ids.update(_delete_repeated_components(slide, repeated_components, protected_ids, logger=logger))

    for item in shapes:
        if item["id"] in deleted_ids:
            continue
        if item["id"] in filled_component_ids:
            continue
        if item["id"] in protected_ids:
            continue
        if item is title_shape or item is active_body_shape:
            continue

        if _looks_like_template_text(item["text"]):
            deleted_ids.update(_delete_unused_component(slide, item, protected_ids=protected_ids, logger=logger))
        elif item["id"] in style.deletable_text_ids or _should_delete_unmapped_content_text(slide, item):
            if logger:
                logger.write(f"delete unmapped content text shape {_item_summary(item)}")
            _delete_shape(item["shape"])
            deleted_ids.add(item["id"])
        elif logger:
            logger.write(f"preserve unmapped text shape {_item_summary(item)}")

    _fit_generated_text_shapes(slide, generated_text_ids, logger=logger)
    _normalize_font_size_for_ids(slide, repeated_heading_ids, "repeated heading", logger=logger)
    _normalize_font_size_for_ids(slide, repeated_description_ids, "repeated description", logger=logger)
    _fit_generated_text_shapes(slide, generated_text_ids, logger=logger)


# Writes speaker notes after COM save using python-pptx.
def _write_notes_with_python_pptx(output_file: Path, plan) -> None:
    prs = Presentation(output_file)

    if len(prs.slides) >= 1:
        prs.slides[0].notes_slide.notes_text_frame.text = ""

    for index, slide_data in enumerate(plan["slides"], start=1):
        if index >= len(prs.slides):
            break
        prs.slides[index].notes_slide.notes_text_frame.text = slide_data.get(
            "speaker_notes",
            "",
        )

    prs.save(output_file)


# Builds a presentation by duplicating detected sample slides through PowerPoint COM.
def build_presentation_from_template_com(
    plan,
    output_file,
    profile: TemplateLayoutProfile,
    debug_log_path=None,
):
    try:
        import pythoncom
        import win32com.client as win32
    except ImportError as exc:
        raise RuntimeError("PowerPoint COM template building requires pywin32.") from exc

    if not profile.title_slide_index or not profile.content_slide_index:
        raise ValueError("COM template builder requires detected title and content sample slides.")

    output_file = Path(output_file).resolve()
    output_file.parent.mkdir(parents=True, exist_ok=True)
    debug_log_path = (
        Path(debug_log_path)
        if debug_log_path
        else output_file.with_name(f"{output_file.stem}_template_debug.log")
    )
    logger = _DebugLogger(debug_log_path)
    logger.write("starting COM template build")
    logger.write(f"template_path={profile.template_path.resolve()}")
    logger.write(f"output_file={output_file}")
    logger.write(
        "profile: "
        f"title_slide={profile.title_slide_index}, content_slide={profile.content_slide_index}, "
        f"confidence={profile.confidence:.2f}, reason={profile.selection_reason}"
    )
    logger.write(f"plan title='{_preview(plan.get('title', ''))}'")
    logger.write(f"plan content slide count={len(plan.get('slides', []))}")

    pythoncom.CoInitialize()
    powerpoint = win32.DispatchEx("PowerPoint.Application")
    try:
        powerpoint.DisplayAlerts = 0
    except Exception:
        pass

    presentation = None
    try:
        presentation = powerpoint.Presentations.Open(str(profile.template_path.resolve()))
        logger.write(f"opened template in PowerPoint with {presentation.Slides.Count} slide(s)")
        original_slide_ids = [
            presentation.Slides(index).SlideID
            for index in range(1, presentation.Slides.Count + 1)
        ]
        title_source_id = presentation.Slides(profile.title_slide_index).SlideID
        content_source_id = presentation.Slides(profile.content_slide_index).SlideID
        logger.write(
            f"source ids: title_source_id={title_source_id}, "
            f"content_source_id={content_source_id}, original_slide_ids={original_slide_ids}"
        )

        keep_ids = []

        title_slide = _duplicate_to_end(presentation, title_source_id)
        keep_ids.append(title_slide.SlideID)
        logger.write(
            f"duplicated title sample slide {profile.title_slide_index} "
            f"to generated slide id={title_slide.SlideID}"
        )
        _replace_title_slide_text(title_slide, plan["title"], logger=logger)

        for slide_index, slide_data in enumerate(plan["slides"], start=1):
            content_slide = _duplicate_to_end(presentation, content_source_id)
            keep_ids.append(content_slide.SlideID)
            logger.write(
                f"duplicated content sample slide {profile.content_slide_index} "
                f"to generated content slide {slide_index} id={content_slide.SlideID}"
            )
            _replace_content_slide_text(
                content_slide,
                slide_data["title"],
                slide_data.get("bullets", []),
                logger=logger,
            )

        for slide_id in reversed(original_slide_ids):
            logger.write(f"deleting original template slide id={slide_id}")
            _slide_by_id(presentation, slide_id).Delete()

        logger.write(f"saving generated presentation to {output_file}")
        presentation.SaveAs(str(output_file))
    except Exception as exc:
        logger.write(f"COM template build failed: {type(exc).__name__}: {exc}")
        raise
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        try:
            powerpoint.Quit()
        except Exception:
            pass
        pythoncom.CoUninitialize()

    logger.write("writing speaker notes with python-pptx")
    _write_notes_with_python_pptx(output_file, plan)
    logger.write("COM template build complete")
    return output_file
