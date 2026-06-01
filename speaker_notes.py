from pathlib import Path
import sys

from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation

load_dotenv()
client = OpenAI()

SYSTEM_PROMPT = """
You are creating speaker notes for a professional business presentation.
Write natural presenter notes, not slide text.
Keep it conversational, crisp, and suitable for narration.
Do not mention "this slide says".
"""


def load_context(path):
    return Path(path).read_text(encoding="utf-8")


def extract_slide_text(slide):
    texts = []
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text.strip():
            texts.append(shape.text.strip())
    return "\n".join(texts)


def generate_notes(slide_number, slide_text, presentation_context):
    prompt = f"""
Overall presentation context:
{presentation_context}

Create speaker notes for slide {slide_number}.

Slide content:
{slide_text}

Output:
- 90 to 130 words
- Presenter-friendly
- Clear transitions
- No bullet list unless necessary
"""

    response = client.responses.create(
        model="gpt-4.1-mini",
        input=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": prompt},
        ],
    )

    return response.output_text.strip()


def run_speaker_notes(
    input_path,
    context_path=None,
    context_text=None,
    output_path=None,
    progress_callback=None,
):
    input_path = Path(input_path)

    if context_text is not None:
        presentation_context = context_text
    elif context_path:
        presentation_context = load_context(context_path)
    else:
        presentation_context = "General business presentation."

    output_path = (
        Path(output_path)
        if output_path
        else input_path.with_name(f"{input_path.stem}_with_ai_notes.pptx")
    )

    ppt = Presentation(input_path)
    slides_to_process = []

    for idx, slide in enumerate(ppt.slides, start=1):
        slide_text = extract_slide_text(slide)
        if slide_text:
            slides_to_process.append((idx, slide, slide_text))

    total = len(slides_to_process)
    if total == 0:
        raise ValueError("No slide text found in the presentation.")

    for step, (idx, slide, slide_text) in enumerate(slides_to_process, start=1):
        if progress_callback:
            progress_callback(
                step,
                total,
                f"Generating speaker notes for slide {idx} ({step}/{total})",
            )

        notes = generate_notes(idx, slide_text, presentation_context)
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    ppt.save(output_path)
    return output_path


def main():
    if len(sys.argv) < 2:
        print("Usage:")
        print("python speaker_notes.py your_presentation.pptx [presentation_context.txt]")
        sys.exit(1)

    input_path = sys.argv[1]
    context_path = sys.argv[2] if len(sys.argv) > 2 else None

    output_path = run_speaker_notes(input_path, context_path=context_path)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()
