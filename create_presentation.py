import json
import sys
from pathlib import Path

from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.util import Pt

from template_com_presentation_builder import build_presentation_from_template_com
from template_layout_detector import (
    detect_template_layout,
    format_template_detection_report,
)

load_dotenv()
client = OpenAI()

SYSTEM_PROMPT = """
You create professional PowerPoint presentations.
Return ONLY valid JSON.
Do not include markdown.
"""


def safe_print(text: str) -> None:
    encoding = sys.stdout.encoding or "utf-8"
    safe_text = text.encode(encoding, errors="replace").decode(encoding)
    print(safe_text, flush=True)


def load_file(path):
    return Path(path).read_text(encoding="utf-8")


def create_slide_plan(content, topic="Business Presentation", num_slides=8):
    content_slide_count = max(1, num_slides - 1)

    prompt = f"""
Create a PowerPoint slide plan from the content below.

Topic:
{topic}

Content:
{content}

Return JSON in this exact format:
{{
  "title": "Presentation title",
  "slides": [
    {{
      "title": "Slide title",
      "bullets": ["bullet 1", "bullet 2", "bullet 3"],
      "speaker_notes": "Natural presenter notes for this slide."
    }}
  ]
}}

Rules:
- Exactly {content_slide_count} content slides in the "slides" array
- 3 to 5 bullets per slide
- Executive-friendly
- Clear storyline
- Speaker notes should be 80 to 130 words
"""

    response = client.responses.create(
        model="gpt-4.1-mini",
        input=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": prompt},
        ],
    )

    return json.loads(response.output_text)


def add_title_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generated with OpenAI"


def add_content_slide(prs, slide_data):
    slide = prs.slides.add_slide(prs.slide_layouts[1])

    slide.shapes.title.text = slide_data["title"]

    body = slide.placeholders[1].text_frame
    body.clear()

    for index, bullet in enumerate(slide_data["bullets"]):
        p = body.paragraphs[0] if index == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(22)

    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = slide_data.get("speaker_notes", "")


def build_presentation(plan, output_file):
    prs = Presentation()

    add_title_slide(prs, plan["title"])

    for slide_data in plan["slides"]:
        add_content_slide(prs, slide_data)

    prs.save(output_file)


def build_presentation_with_optional_template(
    plan,
    output_file,
    template_path=None,
    logs_dir=None,
    progress_callback=None,
):
    if not template_path:
        build_presentation(plan, output_file)
        return output_file

    try:
        profile = detect_template_layout(template_path)
        output_file = Path(output_file)
        logs_dir = Path(logs_dir) if logs_dir else output_file.with_name("logs")
        logs_dir.mkdir(parents=True, exist_ok=True)
        detection_report_path = logs_dir / f"{output_file.stem}_template_detection.log"
        detection_report = format_template_detection_report(template_path, profile)
        detection_report_path.write_text(detection_report, encoding="utf-8")
        safe_print(detection_report)

        if not profile.title_slide_index or not profile.content_slide_index:
            raise ValueError("Could not identify title and content sample slides in template.")

        debug_log_path = logs_dir / f"{output_file.stem}_template_debug.log"
        return build_presentation_from_template_com(
            plan,
            output_file,
            profile,
            debug_log_path=debug_log_path,
        )
    except Exception as exc:
        output_file = Path(output_file)
        logs_dir = Path(logs_dir) if logs_dir else output_file.with_name("logs")
        logs_dir.mkdir(parents=True, exist_ok=True)
        fallback_log_path = logs_dir / f"{output_file.stem}_template_fallback.log"
        fallback_message = (
            f"Template styling failed with {type(exc).__name__}: {exc}\n"
            "Falling back to default PowerPoint layout.\n"
        )
        fallback_log_path.write_text(fallback_message, encoding="utf-8")
        safe_print(fallback_message)
        if progress_callback:
            progress_callback(
                2,
                3,
                f"Template styling failed ({exc}); using default PowerPoint layout",
            )
        build_presentation(plan, output_file)
        return output_file


def run_create_presentation(
    content_path,
    output_path=None,
    topic="Business Presentation",
    num_slides=8,
    template_path=None,
    logs_dir=None,
    progress_callback=None,
):
    content_path = Path(content_path)
    content = load_file(content_path)

    output_path = (
        Path(output_path)
        if output_path
        else content_path.with_suffix(".pptx")
    )

    if progress_callback:
        progress_callback(1, 3, f"Creating slide plan with OpenAI ({num_slides} slides)")

    plan = create_slide_plan(content, topic=topic, num_slides=num_slides)

    if progress_callback:
        if template_path:
            progress_callback(2, 3, "Building presentation from template")
        else:
            progress_callback(2, 3, f"Building presentation ({len(plan['slides'])} slides)")

    build_presentation_with_optional_template(
        plan=plan,
        output_file=output_path,
        template_path=template_path,
        logs_dir=logs_dir,
        progress_callback=progress_callback,
    )

    if progress_callback:
        progress_callback(3, 3, "Presentation saved")

    return output_path


def main():
    if len(sys.argv) < 2:
        print("Usage:")
        print("python create_presentation.py content.txt [num_slides]")
        sys.exit(1)

    content_file = sys.argv[1]
    num_slides = int(sys.argv[2]) if len(sys.argv) > 2 else 8

    if num_slides < 2:
        print("Number of slides must be at least 2 (title slide + 1 content slide).")
        sys.exit(1)

    output_path = run_create_presentation(content_file, num_slides=num_slides)
    print(f"Saved: {output_path}")


if __name__ == "__main__":
    main()
