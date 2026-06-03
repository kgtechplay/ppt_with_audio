from __future__ import annotations

from dataclasses import dataclass
import json
import os
from pathlib import Path

from pptx import Presentation


# Stores the selected title/content sample slides and detector confidence for a template.
@dataclass
class TemplateStyleProfile:
    template_path: Path
    title_slide_index: int | None
    content_slide_index: int | None
    title_score: float
    content_score: float
    confidence: float
    selection_reason: str


# Keeps older imports working while the detector is now sample-slide based.
TemplateLayoutProfile = TemplateStyleProfile


# Builds a compact debug summary for a shape on a sampled slide.
def _placeholder_summary(shape) -> dict:
    summary = {
        "name": shape.name,
        "is_placeholder": bool(shape.is_placeholder),
        "left": int(shape.left),
        "top": int(shape.top),
        "width": int(shape.width),
        "height": int(shape.height),
    }
    if shape.is_placeholder:
        summary["idx"] = shape.placeholder_format.idx
        summary["type"] = str(shape.placeholder_format.type)
    if getattr(shape, "has_text_frame", False):
        summary["text"] = shape.text_frame.text[:120]
    return summary


# Returns non-empty text shapes from a slide.
def _slide_text_shapes(slide) -> list:
    return [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
    ]


# Computes simple text-density and position features for slide classification.
def _slide_features(slide) -> dict:
    text_shapes = _slide_text_shapes(slide)
    texts = [shape.text_frame.text.strip() for shape in text_shapes]
    line_count = sum(len([line for line in text.splitlines() if line.strip()]) for text in texts)
    bulletish_count = sum(text.count("\n") for text in texts)
    largest_area = max((shape.width * shape.height for shape in text_shapes), default=0)
    top_text_count = sum(1 for shape in text_shapes if shape.top < 2_000_000)

    return {
        "text_shape_count": len(text_shapes),
        "line_count": line_count,
        "bulletish_count": bulletish_count,
        "largest_area": largest_area,
        "top_text_count": top_text_count,
        "texts": texts,
    }


# Scores how strongly a slide resembles a title/opening slide.
def _score_title_slide(features: dict) -> float:
    score = 0.0
    score += 3.0 if features["text_shape_count"] <= 3 else 0.0
    score += 2.0 if 1 <= features["line_count"] <= 4 else 0.0
    score += 1.0 if features["top_text_count"] else 0.0
    score -= 2.0 if features["line_count"] > 8 else 0.0
    score -= 1.0 if features["bulletish_count"] > 3 else 0.0
    return score


# Scores how strongly a slide resembles a content/body slide.
def _score_content_slide(features: dict) -> float:
    score = 0.0
    score += 2.5 if features["text_shape_count"] >= 2 else 0.0
    score += 3.0 if features["line_count"] >= 5 else 0.0
    score += 1.0 if features["bulletish_count"] >= 2 else 0.0
    score += 0.5 if features["top_text_count"] else 0.0
    score -= 1.5 if features["line_count"] <= 3 else 0.0
    return score


# Scores every slide once so analysis and reporting share the same data.
def _score_template_slides(prs: Presentation) -> list[dict]:
    summaries = []
    for index, slide in enumerate(prs.slides, start=1):
        features = _slide_features(slide)
        summaries.append(
            {
                "slide_index": index,
                "title_score": _score_title_slide(features),
                "content_score": _score_content_slide(features),
                "features": features,
                "text_shapes": [
                    _placeholder_summary(shape)
                    for shape in _slide_text_shapes(slide)
                ],
            }
        )
    return summaries


# Converts scored slides into the compact payload sent to the AI selector.
def _ai_slide_payload(slide_summaries: list[dict]) -> list[dict]:
    payload = []
    for item in slide_summaries:
        preview = " | ".join(item["features"]["texts"])[:500].replace("\n", " ")
        payload.append(
            {
                "slide_index": item["slide_index"],
                "title_score": item["title_score"],
                "content_score": item["content_score"],
                "text_shape_count": item["features"]["text_shape_count"],
                "line_count": item["features"]["line_count"],
                "bulletish_count": item["features"]["bulletish_count"],
                "top_text_count": item["features"]["top_text_count"],
                "text_preview": preview,
            }
        )
    return payload


# Parses JSON from an AI response, including fenced or prefixed JSON text.
def _parse_ai_json(text: str) -> dict:
    cleaned = (text or "").strip()
    if not cleaned:
        raise ValueError("AI returned an empty response.")

    if cleaned.startswith("```"):
        lines = cleaned.splitlines()
        if lines and lines[0].startswith("```"):
            lines = lines[1:]
        if lines and lines[-1].startswith("```"):
            lines = lines[:-1]
        cleaned = "\n".join(lines).strip()

    try:
        return json.loads(cleaned)
    except json.JSONDecodeError:
        start = cleaned.find("{")
        end = cleaned.rfind("}")
        if start >= 0 and end > start:
            try:
                return json.loads(cleaned[start:end + 1])
            except json.JSONDecodeError as exc:
                raise ValueError(
                    f"AI returned malformed JSON: {cleaned[:300]}"
                ) from exc
        raise ValueError(f"AI returned non-JSON text: {cleaned[:300]}")


# Builds a profile from explicit title/content slide choices.
def _profile_from_slide_choices(
    template_path: Path,
    slide_summaries: list[dict],
    title_slide_index: int,
    content_slide_index: int,
    confidence: float | None = None,
    reason: str | None = None,
) -> TemplateStyleProfile:
    by_index = {item["slide_index"]: item for item in slide_summaries}
    if title_slide_index not in by_index:
        raise ValueError(f"AI selected unknown title slide {title_slide_index}.")
    if content_slide_index not in by_index:
        raise ValueError(f"AI selected unknown content slide {content_slide_index}.")
    if len(slide_summaries) > 1 and title_slide_index == content_slide_index:
        raise ValueError("AI selected the same slide for title and content samples.")

    title_score = max(by_index[title_slide_index]["title_score"], 0.0)
    content_score = max(by_index[content_slide_index]["content_score"], 0.0)
    heuristic_confidence = ((title_score / 6.0) + (content_score / 7.0)) / 2
    selected_confidence = heuristic_confidence if confidence is None else confidence
    selected_confidence = min(0.95, max(0.15, selected_confidence))

    selection_reason = reason or (
        f"Selected slide {title_slide_index} as the title visual sample "
        f"and slide {content_slide_index} as the content visual sample."
    )

    return TemplateStyleProfile(
        template_path=template_path,
        title_slide_index=title_slide_index,
        content_slide_index=content_slide_index,
        title_score=title_score,
        content_score=content_score,
        confidence=selected_confidence,
        selection_reason=selection_reason,
    )


# Selects title/content sample slides with deterministic local scoring.
def _select_template_slides_heuristic(
    template_path: Path,
    slide_summaries: list[dict],
) -> TemplateStyleProfile:
    title_slide = max(slide_summaries, key=lambda item: item["title_score"])
    content_candidates = [
        item for item in slide_summaries if item["slide_index"] != title_slide["slide_index"]
    ] or slide_summaries
    content_slide = max(content_candidates, key=lambda item: item["content_score"])
    return _profile_from_slide_choices(
        template_path,
        slide_summaries,
        title_slide["slide_index"],
        content_slide["slide_index"],
    )


# Returns the highest-scoring slide candidates for a role.
def _top_candidates(slide_summaries: list[dict], score_key: str, exclude_index: int | None = None) -> list[dict]:
    candidates = [
        item for item in slide_summaries
        if exclude_index is None or item["slide_index"] != exclude_index
    ] or slide_summaries
    top_score = max(item[score_key] for item in candidates)
    return [
        item for item in candidates
        if abs(item[score_key] - top_score) <= 0.25
    ]


# Identifies whether heuristic title/content choices are ambiguous enough to ask AI.
def _ai_detection_targets(slide_summaries: list[dict], profile: TemplateStyleProfile) -> tuple[list[str], list[str]]:
    target_roles = []
    reasons = []

    title_ties = _top_candidates(slide_summaries, "title_score")
    if len(title_ties) > 1:
        target_roles.append("title")
        reasons.append(
            "multiple slides tied for best title score: "
            + ", ".join(str(item["slide_index"]) for item in title_ties)
        )
    elif profile.title_score < 3.0:
        target_roles.append("title")
        reasons.append(f"title confidence is weak from scoring alone: {profile.title_score:.1f}")

    content_ties = _top_candidates(
        slide_summaries,
        "content_score",
        exclude_index=profile.title_slide_index,
    )
    if len(content_ties) > 1:
        target_roles.append("content")
        reasons.append(
            "multiple slides tied for best content score: "
            + ", ".join(str(item["slide_index"]) for item in content_ties)
        )
    elif profile.content_score < 4.0:
        target_roles.append("content")
        reasons.append(f"content confidence is weak from scoring alone: {profile.content_score:.1f}")

    return target_roles, reasons


# Selects title/content sample slides with an LLM using the same scored slide summaries.
def _select_template_slides_ai(
    template_path: Path,
    slide_summaries: list[dict],
    target_roles: list[str] | None = None,
    heuristic_profile: TemplateStyleProfile | None = None,
    ambiguity_reasons: list[str] | None = None,
) -> TemplateStyleProfile:
    from openai import OpenAI

    heuristic_profile = heuristic_profile or _select_template_slides_heuristic(template_path, slide_summaries)
    target_roles = target_roles or ["title", "content"]
    ambiguity_reasons = ambiguity_reasons or []

    client = OpenAI()
    model = os.getenv("OPENAI_TEMPLATE_DETECTOR_MODEL", "gpt-4.1-mini")
    prompt = {
        "task": (
            "Choose reusable actual PowerPoint sample slides to duplicate. Resolve only the "
            "roles listed in target_roles. Keep the heuristic choice for any role not listed."
        ),
        "target_roles": target_roles,
        "ambiguity_reasons": ambiguity_reasons,
        "heuristic_choice": {
            "title_slide_index": heuristic_profile.title_slide_index,
            "content_slide_index": heuristic_profile.content_slide_index,
            "title_score": heuristic_profile.title_score,
            "content_score": heuristic_profile.content_score,
            "confidence": heuristic_profile.confidence,
        },
        "rules": [
            "Return JSON only.",
            "Use existing slide_index values.",
            "Use different slides when more than one slide exists.",
            "For target role 'title', pick the best title/cover/opening sample slide.",
            "For target role 'content', pick the best reusable content/control slide, not necessarily the first tied slide.",
            "Prefer content slide with repeated components that can map cleanly to generated bullets: cards, feature blocks, product modules, rows, or labelled content areas. Avoid process/timeline slides when they contain many step boxes that are too specific to the original story.",
            "A good content/control slide usually has reusable structure: repeated rows/cards, body regions, comparison/table/card layouts, or clear content placeholders for title and descriptions.",
            "If a role is not in target_roles, return the heuristic_choice value for that role.",
            "confidence must be a number from 0.0 to 1.0.",
        ],
        "slides": _ai_slide_payload(slide_summaries),
        "response_schema": {
            "title_slide_index": "integer",
            "content_slide_index": "integer",
            "confidence": "number",
            "reason": "short string",
        },
    }
    response = client.responses.create(
        model=model,
        input=[
            {
                "role": "system",
                "content": (
                    "You identify reusable sample slides in PowerPoint templates. "
                    "Return only a single JSON object with no markdown."
                ),
            },
            {"role": "user", "content": json.dumps(prompt)},
        ],
    )
    data = _parse_ai_json(response.output_text)
    title_slide_index = int(data.get("title_slide_index", heuristic_profile.title_slide_index))
    content_slide_index = int(data.get("content_slide_index", heuristic_profile.content_slide_index))
    return _profile_from_slide_choices(
        template_path,
        slide_summaries,
        title_slide_index,
        content_slide_index,
        float(data.get("confidence", 0.5)),
        (
            f"AI selected sample slides for {', '.join(target_roles)}: "
            f"{data.get('reason', '').strip()}"
        ),
    )


# Detects the template sample slides, optionally using AI before heuristic fallback.
def analyze_template_style(template_path, use_ai: bool | str = "auto") -> TemplateStyleProfile:
    template_path = Path(template_path)
    prs = Presentation(template_path)
    slide_summaries = _score_template_slides(prs)

    if not slide_summaries:
        raise ValueError("Template does not contain any slides to use as visual samples.")

    heuristic_profile = _select_template_slides_heuristic(template_path, slide_summaries)
    target_roles, ambiguity_reasons = _ai_detection_targets(slide_summaries, heuristic_profile)
    should_use_ai = use_ai is True or use_ai == "always" or (use_ai == "auto" and target_roles)

    if should_use_ai:
        try:
            return _select_template_slides_ai(
                template_path,
                slide_summaries,
                target_roles=target_roles or ["title", "content"],
                heuristic_profile=heuristic_profile,
                ambiguity_reasons=ambiguity_reasons,
            )
        except Exception as exc:
            heuristic_profile.selection_reason += (
                f" AI selection failed with {type(exc).__name__}: {exc}; "
                "used heuristic slide scoring."
            )
            return heuristic_profile

    return heuristic_profile


# Public wrapper for template sample-slide detection.
def detect_template_layout(template_path, use_ai: bool | str = "auto") -> TemplateLayoutProfile:
    profile = analyze_template_style(template_path, use_ai=use_ai)
    return profile


# Produces per-slide scoring details for logs and debugging.
def summarize_template_slides(template_path) -> list[dict]:
    prs = Presentation(template_path)
    return _score_template_slides(prs)


# Formats the selected samples and all slide scores as a readable report.
def format_template_detection_report(template_path, profile: TemplateStyleProfile) -> str:
    lines = [
        "Template detection report",
        f"Template: {Path(template_path)}",
        "",
        "Selected sample slides",
        f"- Title sample slide: {profile.title_slide_index or 'not needed'}",
        f"- Content sample slide: {profile.content_slide_index or 'not needed'}",
        f"- Confidence: {profile.confidence:.2f}",
        f"- Reason: {profile.selection_reason}",
        "",
        "All slide scores",
    ]

    for item in summarize_template_slides(template_path):
        marker = []
        if item["slide_index"] == profile.title_slide_index:
            marker.append("TITLE_SAMPLE")
        if item["slide_index"] == profile.content_slide_index:
            marker.append("CONTENT_SAMPLE")
        marker_text = f" [{' / '.join(marker)}]" if marker else ""
        preview = " | ".join(item["features"]["texts"])[:180].replace("\n", " ")
        lines.append(
            f"- slide {item['slide_index']}{marker_text} | "
            f"title_score={item['title_score']:.1f}, "
            f"content_score={item['content_score']:.1f}, "
            f"text_shapes={item['features']['text_shape_count']}, "
            f"lines={item['features']['line_count']} | {preview}"
        )

    return "\n".join(lines)
