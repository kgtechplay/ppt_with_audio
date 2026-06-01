import os
import argparse
from pathlib import Path

import requests
from dotenv import load_dotenv
from pptx import Presentation

import win32com.client as win32


load_dotenv()


def extract_speaker_notes(pptx_path):
    prs = Presentation(pptx_path)
    notes_by_slide = {}

    for idx, slide in enumerate(prs.slides, start=1):
        if not slide.has_notes_slide:
            continue

        notes_frame = slide.notes_slide.notes_text_frame
        notes_text = "\n".join(
            p.text.strip()
            for p in notes_frame.paragraphs
            if p.text and p.text.strip()
        )

        if notes_text:
            notes_by_slide[idx] = notes_text

    return notes_by_slide


def generate_elevenlabs_audio(
    text,
    output_file,
    api_key,
    voice_id,
    model_id="eleven_multilingual_v2",
):
    url = (
        f"https://api.elevenlabs.io/v1/text-to-speech/"
        f"{voice_id}?output_format=mp3_44100_128"
    )

    headers = {
        "xi-api-key": api_key,
        "Content-Type": "application/json",
    }

    payload = {
        "text": text,
        "model_id": model_id,
        "voice_settings": {
            "stability": 0.55,
            "similarity_boost": 0.75,
        },
    }

    response = requests.post(url, headers=headers, json=payload, timeout=120)

    if response.status_code != 200:
        raise RuntimeError(
            f"ElevenLabs error {response.status_code}: {response.text}"
        )

    output_file.write_bytes(response.content)


def generate_audio_files(
    notes_by_slide,
    audio_dir,
    api_key,
    voice_id,
    force=False,
    progress_callback=None,
):
    audio_dir.mkdir(parents=True, exist_ok=True)
    audio_map = {}
    total = len(notes_by_slide)

    for step, (slide_num, notes_text) in enumerate(notes_by_slide.items(), start=1):
        audio_file = audio_dir / f"slide_{slide_num:03}.mp3"

        if audio_file.exists() and not force:
            if progress_callback:
                progress_callback(
                    step,
                    total,
                    f"Using existing audio for slide {slide_num} ({step}/{total})",
                )
        else:
            if progress_callback:
                progress_callback(
                    step,
                    total,
                    f"Generating audio for slide {slide_num} ({step}/{total})",
                )
            generate_elevenlabs_audio(
                text=notes_text,
                output_file=audio_file,
                api_key=api_key,
                voice_id=voice_id,
            )

        audio_map[slide_num] = audio_file

    return audio_map


def add_audio_to_ppt(
    input_pptx,
    output_pptx,
    audio_map,
    show_powerpoint=False,
    progress_callback=None,
):
    powerpoint = win32.Dispatch("PowerPoint.Application")
    # Office 365+ often rejects Visible=False ("Hiding the application window
    # is not allowed"). Only force visible when requested; otherwise leave default.
    if show_powerpoint:
        powerpoint.Visible = True
    else:
        try:
            powerpoint.DisplayAlerts = 0  # ppAlertsNone — fewer dialogs during automation
        except Exception:
            pass

    presentation = powerpoint.Presentations.Open(str(Path(input_pptx).resolve()))
    total = len(audio_map)

    for step, (slide_num, audio_file) in enumerate(audio_map.items(), start=1):
        if progress_callback:
            progress_callback(
                step,
                total,
                f"Adding audio to slide {slide_num} ({step}/{total})",
            )

        slide = presentation.Slides(slide_num)

        audio_shape = slide.Shapes.AddMediaObject2(
            FileName=str(audio_file.resolve()),
            LinkToFile=False,
            SaveWithDocument=True,
            Left=20,
            Top=20,
            Width=40,
            Height=40,
        )

        audio_shape.AnimationSettings.PlaySettings.PlayOnEntry = True
        audio_shape.AnimationSettings.PlaySettings.HideWhileNotPlaying = True

    presentation.SaveAs(str(Path(output_pptx).resolve()))
    presentation.Close()
    powerpoint.Quit()


def run_elevenlabs_audio(
    input_pptx,
    audio_dir=None,
    voice_id=None,
    force=False,
    progress_callback=None,
):
    """Generate slide_XXX.mp3 files from speaker notes (no PowerPoint merge)."""
    api_key = os.getenv("ELEVENLABS_API_KEY")
    voice_id = voice_id or os.getenv("ELEVENLABS_VOICE_ID")

    if not api_key:
        raise ValueError("Missing ELEVENLABS_API_KEY in .env")

    if not voice_id:
        raise ValueError("Missing ElevenLabs voice ID")

    input_pptx = Path(input_pptx)
    if not input_pptx.exists():
        raise FileNotFoundError(f"Input PPTX not found: {input_pptx}")

    if audio_dir is None:
        audio_dir = input_pptx.parent / "Audio"
    else:
        audio_dir = Path(audio_dir)

    notes_by_slide = extract_speaker_notes(input_pptx)
    if not notes_by_slide:
        raise ValueError("No speaker notes found in the presentation.")

    audio_map = generate_audio_files(
        notes_by_slide=notes_by_slide,
        audio_dir=audio_dir,
        api_key=api_key,
        voice_id=voice_id,
        force=force,
        progress_callback=progress_callback,
    )

    return audio_dir, audio_map


def run_elevenlabs_voice_to_ppt(
    input_pptx,
    output_pptx=None,
    voice_id=None,
    audio_dir=None,
    force=False,
    show_powerpoint=False,
    progress_callback=None,
    embed_audio=True,
):
    """Generate audio from speaker notes; optionally embed into PPT via win32 COM."""
    input_pptx = Path(input_pptx)

    output_pptx = (
        Path(output_pptx)
        if output_pptx
        else input_pptx.with_name(f"{input_pptx.stem}_with_audio.pptx")
    )

    audio_dir, audio_map = run_elevenlabs_audio(
        input_pptx=input_pptx,
        audio_dir=audio_dir,
        voice_id=voice_id,
        force=force,
        progress_callback=progress_callback,
    )

    if not embed_audio:
        return output_pptx, list(audio_map.values())

    add_audio_to_ppt(
        input_pptx=input_pptx,
        output_pptx=output_pptx,
        audio_map=audio_map,
        show_powerpoint=show_powerpoint,
        progress_callback=progress_callback,
    )

    return output_pptx, list(audio_map.values())


def main():
    parser = argparse.ArgumentParser()

    parser.add_argument("pptx", help="Input PowerPoint file with speaker notes")
    parser.add_argument(
        "--output",
        help="Output PowerPoint file",
        default=None,
    )
    parser.add_argument(
        "--voice-id",
        help="ElevenLabs voice ID. Defaults to ELEVENLABS_VOICE_ID from .env",
        default=os.getenv("ELEVENLABS_VOICE_ID"),
    )
    parser.add_argument(
        "--audio-dir",
        help="Folder for MP3 files (default: <input-folder>/Audio)",
        default=None,
    )
    parser.add_argument(
        "--force",
        action="store_true",
        help="Regenerate audio even if MP3 files already exist",
    )
    parser.add_argument(
        "--show-powerpoint",
        action="store_true",
        help="Show PowerPoint while embedding audio",
    )
    parser.add_argument(
        "--audio-only",
        action="store_true",
        help="Generate MP3 files only; do not embed into PowerPoint",
    )

    args = parser.parse_args()

    if args.audio_only:
        audio_dir, audio_map = run_elevenlabs_audio(
            input_pptx=args.pptx,
            voice_id=args.voice_id,
            audio_dir=args.audio_dir,
            force=args.force,
        )
        print(f"Saved {len(audio_map)} audio file(s) in: {audio_dir}")
    else:
        output_pptx, _ = run_elevenlabs_voice_to_ppt(
            input_pptx=args.pptx,
            output_pptx=args.output,
            voice_id=args.voice_id,
            audio_dir=args.audio_dir,
            force=args.force,
            show_powerpoint=args.show_powerpoint,
        )
        print(f"Saved narrated PowerPoint: {output_pptx}")


if __name__ == "__main__":
    main()
